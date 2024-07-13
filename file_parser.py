'''------------------------------------------File Praser-------------------------------------------------------
-> Create file parsers for NLP pipeline, that parses various file types, e.g. PDFs, DOCX, PPTs, XLSX, CSV, TSV, etc. and create logical chunks of each sections.
---------------------------------------------------------------------------------------------------------------
'''

import argparse
import pandas as pd
import io
import PyPDF2
import tabula
from pptx import Presentation
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

class PDFParser:
    #Gathering data from the PDF file and performing initialization:
    def __init__(self, file_path):
        self.reader = PyPDF2.PdfReader(file_path)
        self.chunks = []
    
    #Collecting logical chunks from the gathered data:
    def parse_file(self):
        for page in self.reader.pages:
            text = page.extract_text()
            paragraphs = text.split("\n\n")

            # Dealing with paragraphs:
            for paragraph in paragraphs:
                self.chunks.append({"type": "paragraph", "text": paragraph, "order": len(self.chunks)})

            # Dealing with tables:
            tables = tabula.read_pdf(io.StringIO(page.extract_text()), pages=1, lattice=True)
            for table in tables:
                headers = [col for col in table.columns]
                data = table.values.tolist()
                self.chunks.append({"type": "table", "headers": headers, "data": data, "order": len(self.chunks)})

            # Dealing with images:
            images = page.extract_images()
            for image in images:
                self.chunks.append({"type": "image", "image": image, "order": len(self.chunks)})

        # Arranging logical chunks in correct order:
        self.chunks = sorted(self.chunks, key=lambda x: x["order"])
        return self.chunks
        
class DocxParser:
    #Gathering data from the csv file and performing initialization:
    def __init__(self, file_path):
        self.doc = Document(file_path)
        self.chunks =[]
        
    #Collecting logical chunks from the gathered data: 
    def parse_file(self):
        for para in self.doc.paragraphs:
            if para.style.name.startswith('List'):
                    self.chunks.append({'type': 'list', 'data': para.text})
            else:
                    self.chunks.append(para.text)
        for table in self.doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data.append(row_data)
            self.chunks.append(table_data)
        for img in self.doc.inline_shapes:
            self.chunks.append({'type': 'image', 'data': img._inline.docx.blob})
        return self.chunks
    
class PPTParser:
    def __init__(self, file_path):
        self.presentation = Presentation(file_path)
        self.chunks = []

    def parse_file(self):
        for slide in self.presentation.slides:
            for shape in slide.shapes:

                #dealing with paragraphs:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = ''
                    for paragraph in text_frame.paragraphs:
                        text += paragraph.text + '\n'
                    self.chunks.append({"type": "paragraph", "text": text, "order": len(self.chunks)})
                
                #Dealing with tables:
                elif shape.shape_type == 19:  
                    table = shape.table
                    headers = [cell.text for cell in table.rows[0].cells]
                    data = [[cell.text for cell in row.cells] for row in table.rows[1:]]
                    self.chunks.append({"type": "table", "headers": headers, "data": data, "order": len(self.chunks)})
                
                #Dealing with images:
                elif shape.shape_type == 13: 
                    image = shape.image
                    self.chunks.append({"type": "image", "image": image, "order": len(self.chunks)})

        # Arranging logical chunks in correct order:
        self.chunks = sorted(self.chunks, key=lambda x: x["order"])
        return self.chunks
    
class CSVParser:
    #Gathering data from the csv file and performing initialization: 
    def __init__(self, file_path):
        self.dataframe = pd.read_csv(file_path)
        self.chunks = dict()
    
    #Collecting logical chunks from the gathered data:
    def parse_file(self):
        for i in self.dataframe.columns:
            self.chunks[i] = self.dataframe[i].tolist()
        self.chunks["dataframe"] = self.dataframe
        self.chunks["headers"] = self.dataframe.columns
        return self.chunks
    
class TSVParser:
    #Gathering data from the tsv file and performing initialization:
    def __init__(self, file_path):
        self.dataframe = pd.read_csv(file_path, delimiter = "\t")
        self.chunks = dict() 
    
    #Collecting logical chunks from the gathered data:    
    def parse_file(self):
        for i in self.dataframe.columns:
            self.chunks[i] = self.dataframe[i].tolist()
            self.chunks["dataframe"] = self.dataframe
            self.chunks["headers"] = self.dataframe.columns
        return self.chunks

class FileParser:
    def __init__(self):
        self.file_path = ""
        self.file_type = ""

    def display_menu(self):
        print("File Parser Menu")
        print("----------------")
        print("1. Select File")
        print("2. Parse File")
        print("3. Exit")

    def select_file(self):
        self.file_path = input("Enter the file path: ")
        self.file_type = input("Enter the file type ( csv, tsv, pdf, ppt, docx): ")

        if self.file_type not in ["csv", "tsv", "pdf", "ppt", "docx"]:
            print("Error: Unsupported file type")
            return False
        else:
            print("File selected successfully")
            return True

    def parse_file(self):
        
        if self.file_path == "" or self.file_type == "":
            print("Error: Please select a file first")
            return

        print("Checking file type...")
        
        if self.file_type == "csv":
            if not self.file_path.endswith(".csv"):
                print("Error: File type does not match the file extension")
            else:
                print("File type matches the file extension")
                csv = CSVParser(self.file_path)
                logical_chunks = csv.parse_file()
                print("Parsing csv file...")
                print(logical_chunks)
                
        elif self.file_type == "tsv":
            if not self.file_path.endswith(".tsv"):
                print("Error: File type does not match the file extension")
            else:
                print("File type matches the file extension")
                tsv = TSVParser(self.file_path)
                logical_chunks = tsv.parse_file()
                print("Parsing tsv file...")
                print(logical_chunks)
                
        elif self.file_type == "pdf":
            if not self.file_path.endswith(".pdf"):
                print("Error: File type does not match the file extension")
            else:
                print("File type matches the file extension")
                pdf = PDFParser(self.file_path)
                logical_chunks = pdf.parse_file()
                print("Parsing pdf file...")
                print(logical_chunks)
                
        elif self.file_type == "ppt":
            if not self.file_path.endswith(".pptx"):
                print("Error: File type does not match the file extension")
            else:
                print("File type matches the file extension")
                ppt = PPTParser(self.file_path)
                logical_chunks = ppt.parse_file()
                print("Parsing pptx file...")
                print(logical_chunks)
                
        elif self.file_type == "docx":
            if not self.file_path.endswith(".docx"):
                print("Error: File type does not match the file extension")
            else:
                print("File type matches the file extension")
                docx = DocxParser(self.file_path)
                logical_chunks = docx.parse_file()
                print("Parsing docx file...")
                print(logical_chunks)
                
    def run(self):
        while True:
            self.display_menu()
            choice = input("Enter your choice: ")
            if choice == "1":
                self.select_file()
            elif choice == "2":
                self.parse_file()
            elif choice == "3":
                print("Exiting...")
                break
            else:
                print("Invalid choice. Please try again.")

if __name__ == "__main__":
    file_parser = FileParser()
    file_parser.run()